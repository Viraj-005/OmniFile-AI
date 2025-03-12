import os
import re
import json
import random
import pyperclip # type: ignore
import numpy as np
from io import BytesIO
import streamlit as st
import pandas as pd
import google.generativeai as genai
from dotenv import load_dotenv
from PyPDF2 import PdfReader
import docx2txt # type: ignore
from pptx import Presentation # type: ignore
import matplotlib.pyplot as plt # type: ignore
import networkx as nx
from striprtf.striprtf import rtf_to_text # type: ignore
from matplotlib.patches import Circle, Rectangle # type: ignore

# Load environment variables
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

# Configure the Gemini API
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel('gemini-2.0-flash')

# Set up page configuration
st.set_page_config(
    page_title="OmniFile AI",
    page_icon="🤖",
    layout="centered"
)

# Custom CSS for styling
st.markdown("""
<style>
    .main {
        background-color: #f8f9fa;
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    }
    .chat-container {
        border: 1px solid #dee2e6;
        border-radius: 8px;
        margin: 1rem 0;
        background: white;
        box-shadow: 0 1px 3px rgba(0,0,0,0.05);
    }
    .question-section {
        padding: 1rem;
        background-color: #f1f3f5;
        border-bottom: 1px solid #dee2e6;
        font-weight: 500;
        color: #212529;
        display: flex;
        align-items: center;
        gap: 0.5rem;
    }
    .answer-section {
        padding: 1.5rem;
        position: relative;
        line-height: 1.6;
        font-size: 0.95rem;
        color: #495057;
    }
    .viz-container {
        margin: 1.5rem 0;
        border-radius: 8px;
        padding: 1rem;
        background: white;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
    .preview-card {
        padding: 1rem;
        background-color: white;
        border-radius: 8px;
        margin: 1rem 0;
        box-shadow: 0 1px 3px rgba(0,0,0,0.05);
    }
    .error-message {
        color: #dc3545;
        background-color: #f8d7da;
        padding: 0.75rem;
        border-radius: 4px;
        margin: 1rem 0;
        border: 1px solid #f5c6cb;
    }
</style>
""", unsafe_allow_html=True)

def generate_random_colors(n):
    return ["#"+''.join(random.choices('0123456789ABCDEF', k=6)) for _ in range(n)]

# Initialize session state
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
if 'document_content' not in st.session_state:
    st.session_state.document_content = ""
if 'upload_key' not in st.session_state:
    st.session_state.upload_key = 0
if 'file_metadata' not in st.session_state:
    st.session_state.file_metadata = {}

def get_file_icon(file_name):
    icons = {
        ".pdf": "📕", ".docx": "📘", ".xlsx": "📊", ".txt": "📝",
        ".pptx": "📽️", ".csv": "📓", ".rtf": "📄", ".ipynb": "📔",
        ".py": "🐍", ".java": "☕", ".js": "📜", ".jsx": "⚛️", ".go": "🐹"
    }
    return next((v for k, v in icons.items() if file_name.lower().endswith(k)), "📁")

def get_file_metadata(uploaded_file, content):
    metadata = {
        "name": uploaded_file.name,
        "type": uploaded_file.type.split('/')[-1].upper(),
        "icon": get_file_icon(uploaded_file.name),
        "word_count": len(content.split()),
        "pages": 0,
        "size": f"{len(uploaded_file.getvalue()) // 1024} KB"
    }

    if uploaded_file.type == "application/pdf":
        try:
            with BytesIO(uploaded_file.getvalue()) as pdf_stream:
                metadata["pages"] = len(PdfReader(pdf_stream).pages)
        except Exception as e:
            st.markdown(f'<div class="error-message">PDF Error: {str(e)}</div>', unsafe_allow_html=True)
    return metadata

def process_file(uploaded_file):
    try:
        content = uploaded_file.getvalue()
        file_name = uploaded_file.name.lower()
        
        if file_name.endswith('.pdf'):
            with BytesIO(content) as pdf_stream:
                return "\n".join(page.extract_text() or "" for page in PdfReader(pdf_stream).pages)
        
        elif file_name.endswith('.rtf'):
            rtf_text = content.decode('utf-8')
            return rtf_to_text(rtf_text)
        
        elif file_name.endswith('.ipynb'):
            notebook = json.loads(content)
            text_content = []
            for cell in notebook.get('cells', []):
                if cell['cell_type'] == 'code':
                    text_content.append("```python\n" + '\n'.join(cell['source']) + "\n```")
                elif cell['cell_type'] == 'markdown':
                    text_content.append('\n'.join(cell['source']))
            return '\n\n'.join(text_content)
        
        elif file_name.endswith(('.py', '.java', '.js', '.jsx', '.go')):
            return content.decode("utf-8")
        
        elif file_name.endswith('.docx'):
            return docx2txt.process(BytesIO(content))
        
        elif file_name.endswith('.xlsx'):
            return pd.read_excel(BytesIO(content)).to_string()
        
        elif file_name.endswith('.pptx'):
            prs = Presentation(BytesIO(content))
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        
        else:
            try:
                return content.decode("utf-8")
            except:
                return ""
    
    except Exception as e:
        st.markdown(f'<div class="error-message">Processing Error ({uploaded_file.name}): {str(e)}</div>', unsafe_allow_html=True)
        return ""

def generate_visualization(response, viz_type):
    try:
        # Common data parsing functions
        def parse_numerical_data():
            return [float(x) for line in response.split('\n') 
                   for x in re.findall(r"[-+]?\d*\.\d+|\d+", line)][:100]

        def parse_xy_pairs():
            return [tuple(map(float, re.findall(r"[-+]?\d*\.\d+|\d+", line)[:2]))
                   for line in response.split('\n') if any(c.isdigit() for c in line)][:50]

        # Visualization handlers
        if viz_type == 'histogram':
            data = parse_numerical_data()
            if len(data) > 1:
                plt.figure(figsize=(10,6))
                plt.hist(data, bins=10, color=random.choice(generate_random_colors(1)), edgecolor='black')
                plt.xlabel('Values')
                plt.ylabel('Frequency')
                buf = BytesIO()
                plt.savefig(buf, format="png", dpi=120)
                plt.close()
                return buf.getvalue()

        elif viz_type == 'contour':
            x = np.linspace(-3, 3, 100)
            y = np.linspace(-3, 3, 100)
            X, Y = np.meshgrid(x, y)
            Z = np.sin(X) * np.cos(Y)
            plt.figure(figsize=(10,6))
            plt.contourf(X, Y, Z, 20, cmap='RdGy')
            plt.colorbar()
            buf = BytesIO()
            plt.savefig(buf, format="png", dpi=120)
            plt.close()
            return buf.getvalue()

        elif viz_type == 'er_diagram':
            relations = [tuple(map(str.strip, line.split('--', 1)))
                        for line in response.split('\n') if '--' in line][:15]
            if relations:
                G = nx.Graph()
                for ent1, ent2 in relations:
                    G.add_node(ent1[:15], type='entity')
                    G.add_node(ent2[:15], type='entity')
                    G.add_edge(ent1[:15], ent2[:15], type='relationship')
                
                plt.figure(figsize=(12,8))
                pos = nx.spring_layout(G)
                nx.draw_networkx_nodes(G, pos, node_shape='s', 
                                      node_size=2500, node_color='#ffd700')
                nx.draw_networkx_edges(G, pos, style='dashed')
                nx.draw_networkx_labels(G, pos, font_size=10)
                buf = BytesIO()
                plt.savefig(buf, format="png", dpi=120)
                plt.close()
                return buf.getvalue()

        elif viz_type == 'object_diagram':
            objects = [line.split(':') for line in response.split('\n') if ':' in line][:20]
            if objects:
                plt.figure(figsize=(12,8))
                ax = plt.gca()
                for i, (name, attributes) in enumerate(objects):
                    x = i % 4 * 2.5
                    y = i // 4 * -2.5
                    ax.add_patch(Rectangle((x-1, y-1), 2, 1.5, fill=True, color='#a0c8f0'))
                    plt.text(x-0.9, y-0.8, f"{name}\n{attributes.strip()}", 
                            fontsize=8, va='top')
                plt.axis('equal')
                plt.axis('off')
                buf = BytesIO()
                plt.savefig(buf, format="png", dpi=120)
                plt.close()
                return buf.getvalue()

        elif viz_type == 'state_diagram':
            transitions = [re.split(r'\[|\]', line) for line in response.split('\n') 
                          if '->' in line and '[' in line][:20]
            if transitions:
                G = nx.DiGraph()
                for line in transitions:
                    src, label = line[0].split('->')
                    dest = line[-1].strip()
                    G.add_edge(src.strip(), dest, label=label.strip())
                
                plt.figure(figsize=(12,8))
                pos = nx.circular_layout(G)
                nx.draw_networkx_nodes(G, pos, node_shape='s', 
                                      node_size=2500, node_color='#90EE90')
                nx.draw_networkx_edges(G, pos, arrows=True)
                nx.draw_networkx_edge_labels(G, pos, 
                    edge_labels={(u,v): d['label'] for u,v,d in G.edges(data=True)})
                nx.draw_networkx_labels(G, pos)
                buf = BytesIO()
                plt.savefig(buf, format="png", dpi=120)
                plt.close()
                return buf.getvalue()

        elif viz_type == 'flow_chart':
            steps = [line[2:].strip() for line in response.split('\n') 
                    if line.startswith(('1.', '2.', '3.', '- '))][:10]
            if len(steps) > 1:
                fig, ax = plt.subplots(figsize=(12,8))
                y_pos = 0
                for i, step in enumerate(steps):
                    ax.add_patch(Rectangle((0.1, y_pos), 0.8, 0.1, 
                                        color='#4a86e8', alpha=0.3))
                    plt.text(0.5, y_pos+0.05, step, ha='center', va='center')
                    if i < len(steps)-1:
                        plt.arrow(0.5, y_pos-0.05, 0, -0.1, 
                                head_width=0.05, head_length=0.05, fc='k')
                    y_pos -= 0.2
                plt.axis('off')
                buf = BytesIO()
                plt.savefig(buf, format="png", dpi=120)
                plt.close()
                return buf.getvalue()
        
    except Exception as e:
        st.markdown(f'<div class="error-message">Visualization Error: {str(e)}</div>', unsafe_allow_html=True)
        return None

def check_viz_type(question):
    viz_map = {
        'histogram': ['histogram', 'distribution'],
        'contour': ['contour', '3d plot'],
        'er_diagram': ['er diagram', 'entity relationship'],
        'object_diagram': ['object diagram', 'class diagram'],
        'state_diagram': ['state diagram', 'state machine'],
        'flow_chart': ['flow chart', 'process diagram'],
        'bar_chart': ['bar chart', 'bar graph'],
        'pie_chart': ['pie chart'],
        'line_graph': ['line graph', 'trend line'],
        'table': ['table', 'tabular data']
    }
    q = question.lower()
    for viz_type, keywords in viz_map.items():
        if any(kw in q for kw in keywords):
            return viz_type
    return None

def get_response(question, context):
    try:
        prompt = f"""Analyze the document and answer the question professionally.
        **Document Content:** {context}
        **Question:** {question}
        
        Formatting Rules:
        [Include formatting rules from previous code here]
        """
        
        response = model.generate_content(prompt)
        return response.text
        
    except Exception as e:
        st.markdown(f'<div class="error-message">API Error: {str(e)}</div>', unsafe_allow_html=True)
        return "Failed to generate answer. Please try again."

def main():

    col1, col2, col3 = st.columns([1, 6, 1])
    with col2:
        st.image(
            "robot.jpg",
            use_container_width=False,    # Updated parameter
            caption="Intelligent Document Analysis Platform",
            output_format="PNG",  # Force format if needed
            width=300  # Set specific width if needed
        )
    
    st.title("🤖 OmniFile AI")
    
    # Initialize session state for input clearing
    if 'clear_input' not in st.session_state:
        st.session_state.clear_input = False

    # Sidebar
    with st.sidebar:
        st.header("📂 File Manager")
        
        if st.button("🔄 New Session"):
            st.session_state.upload_key += 1
            st.session_state.chat_history = []
            st.session_state.document_content = ""
            st.session_state.file_metadata = {}
            st.session_state.clear_input = True
            st.rerun()
        
        uploaded_files = st.file_uploader(
            "📤 Upload documents",
            type=[
                "pdf", "docx", "xlsx", "txt", "pptx", "csv",
                "rtf", "ipynb", "py", "java", "js", "jsx", "go"
            ],
            accept_multiple_files=True,
            key=f"uploader_{st.session_state.upload_key}"
        )
        
        if uploaded_files:
            with st.spinner("⏳ Processing files..."):
                st.session_state.document_content = ""
                st.session_state.file_metadata = {}
                
                for file in uploaded_files:
                    content = process_file(file)
                    if content:
                        st.session_state.document_content += f"\n\n--- {file.name} ---\n{content}"
                        st.session_state.file_metadata[file.name] = get_file_metadata(file, content)

                if st.session_state.document_content:
                    st.success(f"✅ Loaded {len(uploaded_files)} file(s)")
                    selected_file = st.selectbox("📄 View File", options=list(st.session_state.file_metadata.keys()))
                    meta = st.session_state.file_metadata[selected_file]
                    st.markdown(f"""
                    <div class="preview-card">
                        <div style="font-size: 2rem">{meta['icon']}</div>
                        <h4>{meta['name']}</h4>
                        <div class="metadata-item">📝 Type: {meta['type']}</div>
                        <div class="metadata-item">🔠 Words: {meta['word_count']}</div>
                        <div class="metadata-item">📦 Size: {meta['size']}</div>
                        {f"<div class='metadata-item'>📑 Pages: {meta['pages']}</div>" if meta['pages'] else ""}
                    </div>
                    """, unsafe_allow_html=True)

    # Main Interface
    st.header("💬 Document Analysis")
    
    # Create form for text input and analyze button
    with st.form(key='analysis_form'):
        question = st.text_input(
            "🔍 Enter your question:",
            placeholder="Ask about the document content...",
            key="question_input",
            value="" if st.session_state.clear_input else st.session_state.get("question_input", "")
        )
        
        # Analyze button inside form
        submitted = st.form_submit_button(
            "🚀 Analyze",
            use_container_width=True,
            type="primary"
        )
    
    # Reset clear flag after rendering input
    if st.session_state.clear_input:
        st.session_state.clear_input = False

    # Handle form submission (either Enter key or button click)
    if submitted:
        if not uploaded_files:
            st.error("⚠️ Please upload documents before analyzing")
        elif not question.strip():
            st.error("⚠️ Please enter a question")
        else:
            if st.session_state.document_content:
                with st.spinner("🔍 Analyzing..."):
                    response = get_response(question, st.session_state.document_content)
                    viz_type = check_viz_type(question)
                    viz = generate_visualization(response, viz_type) if viz_type else None
                    
                    st.session_state.chat_history.insert(0, {
                        "question": question,
                        "answer": response,
                        "viz": viz,
                        "viz_type": viz_type
                    })
                    
                    # Set flag to clear input on next render
                    st.session_state.clear_input = True
                    st.rerun()
            else:
                st.error("⚠️ Failed to process uploaded documents")

    # Display Chat History
    for idx, chat in enumerate(st.session_state.chat_history):
        with st.container():
            st.markdown(f"""
            <div class="chat-container">
                <div class="question-section">🔍 {chat['question']}</div>
                <div class="answer-section">📝 {chat['answer']}</div>
            </div>
            """, unsafe_allow_html=True)
            
            if chat.get('viz') is not None:
                with st.expander("📊 Visualization", expanded=True):
                    if chat['viz_type'] == 'table':
                        st.dataframe(pd.read_csv(BytesIO(chat['viz']), use_container_width=True))
                    else:
                        st.image(chat['viz'], use_container_width=True)

            cols = st.columns([1, 1, 14])
            with cols[0]:
                if st.button("📋",
                    key=f"copy_{idx}_{hash(chat['answer'])}",
                    help="Copy to clipboard"):
                    pyperclip.copy(chat['answer'])
                    st.toast("✅ Copied to clipboard")
            with cols[1]:
                st.download_button(
                    label="⭳",
                    data=chat['answer'],
                    file_name=f"analysis_{idx}.txt",
                    mime="text/plain",
                    key=f"dl_{idx}_{hash(chat['answer'])}",
                    help="Download analysis result"
                )

if __name__ == "__main__":
    main()